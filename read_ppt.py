import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

prs = Presentation(r'c:\Users\Admin\Documents\Student performance predictor\Student_Performance_Prediction_ML_Final.pptx')

for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text.strip())
    print()
